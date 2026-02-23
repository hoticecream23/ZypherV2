"""
Zypher Extractor
Extracts text from documents for search indexing.
Supports PDF, DOCX, XLSX, PPTX, TXT, CSV, Images.
"""
import csv
import os
from pathlib import Path
from typing import Dict, Generator
from ..utils.logger import logger

# Top-level optional imports with availability flags
try:
    import fitz
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import pytesseract
    from PIL import Image
    HAS_OCR = True
except ImportError:
    HAS_OCR = False

try:
    from charset_normalizer import from_path
    HAS_CHARSET = True
except ImportError:
    HAS_CHARSET = False


class Extractor:
    SUPPORTED_FORMATS = {
        '.pdf', '.jpg', '.jpeg', '.png',
        '.tiff', '.docx', '.xlsx', '.pptx', '.txt', '.csv'
    }

    def extract_text_for_search(self, file_path: str) -> Dict:
        """
        Main entry point. Returns status, format, page_count, full_text.
        Consumes the generator internally — use extract_pages_streaming
        for memory-efficient processing of large files.
        """
        ext = Path(file_path).suffix.lower()

        if ext not in self.SUPPORTED_FORMATS:
            return self._error_result(ext, f"Unsupported format: {ext}")

        pages = []
        full_text_parts = []

        try:
            for page in self.extract_pages_streaming(file_path):
                pages.append(page)
                full_text_parts.append(page['text'])

            return {
                'status': 'success',
                'format': ext.lstrip('.'),
                'pages': pages,
                'full_text': '\n'.join(full_text_parts),
                'page_count': len(pages)
            }
        except Exception as e:
            logger.error(f"Extraction failed for {file_path}: {e}", exc_info=True)
            return self._error_result(ext.lstrip('.'), str(e))

    def extract_pages_streaming(self, file_path: str) -> Generator:
        """
        Memory-efficient generator — yields one page/sheet/slide at a time.
        Use this for large files to avoid loading everything into RAM.
        """
        ext = Path(file_path).suffix.lower()

        extractors = {
            '.pdf':  self._stream_pdf,
            '.docx': self._stream_docx,
            '.xlsx': self._stream_xlsx,
            '.pptx': self._stream_pptx,
            '.txt':  self._stream_txt,
            '.csv':  self._stream_csv,
            '.jpg':  self._stream_image,
            '.jpeg': self._stream_image,
            '.png':  self._stream_image,
            '.tiff': self._stream_image,
        }

        if ext not in extractors:
            raise ValueError(f"Unsupported format: {ext}")

        yield from extractors[ext](file_path)

    # -------------------------
    # PDF
    # -------------------------
    def _stream_pdf(self, file_path: str) -> Generator:
        if not HAS_FITZ:
            raise ImportError("PyMuPDF not installed — pip install pymupdf")

        try:
            doc = fitz.open(file_path)
        except fitz.FileDataError as e:
            raise ValueError(f"Corrupted or invalid PDF: {e}")

        try:
            for i in range(len(doc)):
                page = doc[i]
                text = page.get_text("text").strip()

                # Scanned PDF detection — if no text, try OCR
                if not text and HAS_OCR:
                    logger.info(f"Page {i+1} has no text — attempting OCR")
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    text = pytesseract.image_to_string(img).strip()

                yield {
                    'page_num': i + 1,
                    'text': text,
                    'scanned': not page.get_text("text").strip() and bool(text)
                }
        finally:
            doc.close()

    # -------------------------
    # DOCX
    # -------------------------
    def _stream_docx(self, file_path: str) -> Generator:
        if not HAS_DOCX:
            raise ImportError("python-docx not installed — pip install python-docx")

        try:
            doc = DocxDocument(file_path)
        except Exception as e:
            raise ValueError(f"Corrupted or invalid DOCX: {e}")

        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        yield {
            'page_num': 1,
            'text': '\n'.join(paragraphs)
        }

    # -------------------------
    # XLSX
    # -------------------------
    def _stream_xlsx(self, file_path: str) -> Generator:
        if not HAS_OPENPYXL:
            raise ImportError("openpyxl not installed — pip install openpyxl")

        try:
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        except Exception as e:
            raise ValueError(f"Corrupted or invalid XLSX: {e}")

        try:
            for sheet_num, sheet_name in enumerate(wb.sheetnames):
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    row_text = '\t'.join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        rows.append(row_text)
                yield {
                    'page_num': sheet_num + 1,
                    'sheet': sheet_name,
                    'text': '\n'.join(rows)
                }
        finally:
            wb.close()

    # -------------------------
    # PPTX
    # -------------------------
    def _stream_pptx(self, file_path: str) -> Generator:
        if not HAS_PPTX:
            raise ImportError("python-pptx not installed — pip install python-pptx")

        try:
            prs = Presentation(file_path)
        except Exception as e:
            raise ValueError(f"Corrupted or invalid PPTX: {e}")

        for i, slide in enumerate(prs.slides):
            text_parts = []
            self._extract_pptx_shapes(slide.shapes, text_parts)

            # Also extract slide notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    text_parts.append(f"[Notes]: {notes}")

            yield {
                'page_num': i + 1,
                'text': '\n'.join(text_parts)
            }

    def _extract_pptx_shapes(self, shapes, text_parts: list):
        """Recursively extract text from all shape types including tables and groups"""
        for shape in shapes:
            if shape.has_text_frame and shape.text.strip():
                text_parts.append(shape.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = '\t'.join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        text_parts.append(row_text)
            # Recurse into grouped shapes
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                self._extract_pptx_shapes(shape.shapes, text_parts)

    # -------------------------
    # TXT
    # -------------------------
    def _stream_txt(self, file_path: str) -> Generator:
        encoding = self._detect_encoding(file_path)
        try:
            with open(file_path, 'r', encoding=encoding, errors='replace') as f:
                text = f.read().strip()
            yield {'page_num': 1, 'text': text}
        except Exception as e:
            raise ValueError(f"Could not read text file: {e}")

    # -------------------------
    # CSV
    # -------------------------
    def _stream_csv(self, file_path: str) -> Generator:
        """Proper CSV parsing preserving tabular structure"""
        encoding = self._detect_encoding(file_path)
        try:
            rows = []
            with open(file_path, 'r', encoding=encoding, errors='replace', newline='') as f:
                reader = csv.reader(f)
                for row in reader:
                    row_text = '\t'.join(cell.strip() for cell in row if cell.strip())
                    if row_text:
                        rows.append(row_text)
            yield {'page_num': 1, 'text': '\n'.join(rows)}
        except Exception as e:
            raise ValueError(f"Could not read CSV file: {e}")

    # -------------------------
    # Images
    # -------------------------
    def _stream_image(self, file_path: str) -> Generator:
        if not HAS_OCR:
            logger.warning("pytesseract not installed — image OCR skipped")
            yield {'page_num': 1, 'text': ''}
            return

        try:
            text = pytesseract.image_to_string(Image.open(file_path)).strip()
            yield {'page_num': 1, 'text': text}
        except Exception as e:
            raise ValueError(f"Image OCR failed: {e}")

    # -------------------------
    # Helpers
    # -------------------------
    def _detect_encoding(self, file_path: str) -> str:
        """Detect file encoding using charset_normalizer, fallback to utf-8"""
        if HAS_CHARSET:
            try:
                result = from_path(file_path).best()
                if result:
                    return result.encoding
            except Exception:
                pass
        return 'utf-8'

    def _error_result(self, fmt: str, error_message: str) -> Dict:
        return {
            'status': 'error',
            'format': fmt,
            'pages': [],
            'full_text': '',
            'page_count': 0,
            'error_message': error_message
        }


__all__ = ["Extractor"]
'''

Key changes made:
- All imports at top level with availability flags
- Generator-based streaming via `extract_pages_streaming` for memory efficiency
- `status: success/error` and `error_message` on all returns
- Dedicated `_extract_csv` with proper `csv.reader`
- Recursive PPTX shape extraction including tables, groups, and notes
- Encoding detection with `charset_normalizer`
- Scanned PDF detection with OCR fallback
- Specific exception types where possible

One new dependency to install:
'''
